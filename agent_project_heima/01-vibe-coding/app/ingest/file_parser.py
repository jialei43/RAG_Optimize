"""
多格式文件解析器。

支持的格式:
  - .pdf       使用 pypdf 提取文本（不含 OCR）
  - .docx      使用 python-docx，含段落和表格
  - .doc       需要系统安装 antiword 或 LibreOffice，否则跳过
  - .xlsx/.xls 使用 openpyxl + pandas，将表格序列化为文本
  - .csv       使用 pandas，转为文本行
  - .pptx      使用 python-pptx，提取每张幻灯片文本
  - .txt/.md   直接读取 UTF-8 文本
  - 图片        使用 Pillow 获取元信息；无 OCR（需 Tesseract）

解析结果统一返回 str，调用方负责后续分块。
"""

import io  # 字节流（供 Pillow 使用）
import logging  # 记录跳过/警告信息
import os  # 文件后缀判断
import subprocess  # 调用系统命令 antiword
from pathlib import Path  # 路径操作

logger = logging.getLogger(__name__)  # 模块级 logger

# --------------------------------------------------------------------------- #
# 内部解析函数（每种格式独立，便于单独测试）
# --------------------------------------------------------------------------- #

def _parse_pdf(path: str) -> str:
    """从 PDF 提取所有页面文本，每页间空行分隔。"""
    try:
        import pypdf  # noqa: PLC0415  # 惰性导入，避免无 pypdf 时整模块崩溃
        reader = pypdf.PdfReader(path)  # 打开 PDF
        pages = []
        for i, page in enumerate(reader.pages):  # 遍历每页
            text = page.extract_text() or ""  # 提取文本，失败返回空串
            if text.strip():  # 忽略空白页
                pages.append(f"[第{i+1}页]\n{text.strip()}")
        return "\n\n".join(pages)  # 页间双换行分隔
    except ImportError:
        raise RuntimeError("缺少依赖: pypdf，请执行 pip install pypdf")


def _parse_docx(path: str) -> str:
    """从 DOCX 提取段落和表格文本，保留结构换行。"""
    try:
        from docx import Document  # noqa: PLC0415
        doc = Document(path)  # 加载文档
        parts = []

        # 按文档 XML 顺序遍历段落和表格（保持原始顺序）
        for block in doc.element.body:  # type: ignore[attr-defined]
            tag = block.tag.split("}")[-1]  # 去掉命名空间前缀

            if tag == "p":  # 段落
                from docx.oxml.ns import qn  # noqa: PLC0415
                text = "".join(r.text for r in block.iter(qn("w:t")))
                if text.strip():
                    parts.append(text.strip())

            elif tag == "tbl":  # 表格
                from docx.table import Table  # noqa: PLC0415
                table = Table(block, doc)
                rows = []
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    rows.append(" | ".join(cells))  # 单元格用 | 分隔
                if rows:
                    parts.append("\n".join(rows))

        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("缺少依赖: python-docx，请执行 pip install python-docx")


def _parse_doc(path: str) -> str:
    """通过系统命令 antiword 提取 .doc 老格式文本。"""
    if subprocess.run(["which", "antiword"], capture_output=True).returncode != 0:
        # antiword 不存在，尝试 LibreOffice 转换
        if subprocess.run(["which", "libreoffice"], capture_output=True).returncode != 0:
            logger.warning(  # 没有可用工具，跳过该文件
                "跳过 .doc 文件 %s：系统未安装 antiword 或 LibreOffice。"
                "可执行 'brew install antiword' 后重试。",
                path,
            )
            return ""

        # 用 LibreOffice 转为 docx，再递归解析
        import tempfile  # noqa: PLC0415
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx",
                 "--outdir", tmpdir, path],
                capture_output=True, check=True,
            )
            stem = Path(path).stem
            converted = os.path.join(tmpdir, stem + ".docx")
            if os.path.exists(converted):
                return _parse_docx(converted)
        return ""

    result = subprocess.run(["antiword", path], capture_output=True, text=True)
    return result.stdout.strip()


def _parse_excel(path: str) -> str:
    """将 Excel 各 Sheet 数据序列化为可读文本（列名 + 行数据）。"""
    try:
        import pandas as pd  # noqa: PLC0415
        ext = Path(path).suffix.lower()
        engine = "openpyxl" if ext in (".xlsx", ".xlsm") else "xlrd"  # .xls 用 xlrd
        xl = pd.ExcelFile(path, engine=engine)  # 打开工作簿
        parts = []
        for sheet_name in xl.sheet_names:  # 遍历所有 Sheet
            df = xl.parse(sheet_name)
            df = df.fillna("")  # 空单元格替换为空串
            text_lines = [f"[Sheet: {sheet_name}]"]
            text_lines.append(" | ".join(str(c) for c in df.columns))  # 列头
            for _, row in df.iterrows():  # 逐行转文本
                text_lines.append(" | ".join(str(v) for v in row.values))
            parts.append("\n".join(text_lines))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("缺少依赖: openpyxl / pandas，请执行 pip install openpyxl pandas")


def _parse_csv(path: str) -> str:
    """将 CSV 每行序列化为管道符分隔的文本。"""
    try:
        import pandas as pd  # noqa: PLC0415
        df = pd.read_csv(path, encoding="utf-8-sig")  # utf-8-sig 兼容 BOM
        df = df.fillna("")
        lines = [" | ".join(str(c) for c in df.columns)]  # 列头
        for _, row in df.iterrows():
            lines.append(" | ".join(str(v) for v in row.values))
        return "\n".join(lines)
    except UnicodeDecodeError:
        import pandas as pd  # noqa: PLC0415
        df = pd.read_csv(path, encoding="gbk")  # 降级尝试 GBK
        df = df.fillna("")
        lines = [" | ".join(str(c) for c in df.columns)]
        for _, row in df.iterrows():
            lines.append(" | ".join(str(v) for v in row.values))
        return "\n".join(lines)


def _parse_pptx(path: str) -> str:
    """从 PPTX 每张幻灯片提取文本框内容。"""
    try:
        from pptx import Presentation  # noqa: PLC0415
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):  # 遍历幻灯片
            texts = []
            for shape in slide.shapes:  # 遍历形状
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs)
                        if line.strip():
                            texts.append(line.strip())
            if texts:
                slides.append(f"[幻灯片 {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise RuntimeError("缺少依赖: python-pptx，请执行 pip install python-pptx")


def _parse_image(path: str) -> str:
    """返回图片元信息描述；不含 OCR（需系统安装 Tesseract）。"""
    try:
        from PIL import Image  # noqa: PLC0415
        with Image.open(path) as img:
            w, h = img.size
            mode = img.mode  # 颜色模式，如 RGB / RGBA / L
        filename = os.path.basename(path)
        # 无 OCR，仅记录元信息；实际文字内容无法提取
        logger.warning(
            "图片 %s 仅提取元信息，若需 OCR 请安装 tesseract-ocr 并执行 pip install pytesseract",
            path,
        )
        return f"[图片文件: {filename}，尺寸: {w}x{h}，颜色模式: {mode}]"
    except ImportError:
        raise RuntimeError("缺少依赖: Pillow，请执行 pip install Pillow")


def _parse_text(path: str) -> str:
    """读取纯文本文件（.txt / .md）。"""
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb2312"):  # 依次尝试编码
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise RuntimeError(f"无法识别文件 {path} 的编码，请转换为 UTF-8 后重试")


# --------------------------------------------------------------------------- #
# 公共接口
# --------------------------------------------------------------------------- #

# 后缀 → 解析函数映射表
_PARSERS: dict[str, object] = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".doc":  _parse_doc,
    ".xlsx": _parse_excel,
    ".xlsm": _parse_excel,
    ".xls":  _parse_excel,
    ".csv":  _parse_csv,
    ".pptx": _parse_pptx,
    ".txt":  _parse_text,
    ".md":   _parse_text,
    ".png":  _parse_image,
    ".jpg":  _parse_image,
    ".jpeg": _parse_image,
    ".gif":  _parse_image,
    ".bmp":  _parse_image,
    ".webp": _parse_image,
}

SUPPORTED_EXTENSIONS = set(_PARSERS.keys())  # 供调用方做文件过滤


def parse_file(path: str) -> str:
    """
    解析指定文件，返回提取后的纯文本。

    Args:
        path: 文件绝对路径或相对路径。

    Returns:
        提取的文本内容；若格式不支持返回空串并打印警告。

    Raises:
        FileNotFoundError: 文件不存在。
        RuntimeError: 解析过程中遇到依赖缺失等不可恢复错误。
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"文件不存在: {path}")

    ext = Path(path).suffix.lower()  # 统一小写后缀
    parser = _PARSERS.get(ext)

    if parser is None:
        logger.warning("不支持的文件格式 %s，跳过: %s", ext, path)
        return ""

    logger.info("解析文件: %s", path)
    return parser(path)  # type: ignore[call-arg]
